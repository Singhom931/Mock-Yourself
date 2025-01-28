import PyPDF2
from pptx import Presentation
from tkinter import Tk, filedialog

def browseFiles():
    Tk().withdraw()  # Prevent the root window from appearing
    filename = filedialog.askopenfilename(
        initialdir=".", title="Select a File",
        filetypes=(("PDF files", "*.pdf"), ("PPTX files", "*.pptx"), ("all files", "*.*"))
    )
    return filename

def extract_text_from_pdf(pdf_path):
    text = ""
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def main():
    selected_file = browseFiles()
    if selected_file:
        if selected_file.lower().endswith('.pdf'):
            extracted_text = extract_text_from_pdf(selected_file)
            with open("output\extracted.txt", "w", encoding="utf-8") as output_file:
                output_file.write(extracted_text)
            print("Text extracted and saved to 'extracted.txt'")
        elif selected_file.lower().endswith('.pptx'):
            extracted_text = extract_text_from_pptx(selected_file)
            with open("GPT\extracted.txt", "w", encoding="utf-8") as output_file:
                output_file.write(extracted_text)
            print("Text extracted and saved to 'extracted.txt'")
        else:
            print("Unsupported file format.")

if __name__ == '__main__':
    main()
